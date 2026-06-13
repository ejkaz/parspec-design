"""
parspec-pptx — builder helpers.

Low-level primitives every layout uses: new_deck(), blank(), add_text(),
add_rect(), add_line(), corner_brackets(), footer(), motif accents.

Defaults reflect Parspec brand_artifacts mode (dark + Brand Orange).
Layouts in `references/layouts.md` recipe these primitives into slides.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .brand import Brand, load as load_brand


# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Weight semantics for add_text(). python-pptx only exposes bold/not-bold;
# heavier weights (800/900) render only if the install has those Montserrat
# weight files. We mark all heavy weights as bold and rely on font fallback.
WEIGHT_MAP = {
    "regular":   (400, False),
    "medium":    (500, False),
    "semibold":  (600, True),
    "bold":      (700, True),
    "extrabold": (800, True),
    "black":     (900, True),
}


# ── deck + slide ──────────────────────────────────────────────────────
def new_deck() -> Presentation:
    """Return a fresh 16:9 Presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation, *, bg: RGBColor | None = None,
          brand: Brand | None = None):
    """Add a blank slide with Parspec Black (or override) background."""
    if bg is None:
        if brand is None:
            brand = load_brand()
        bg = brand.surface_dark
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.shadow.inherit = False
    return slide


# ── primitives: text, rect, line ──────────────────────────────────────
def add_text(slide, x, y, w, h, text, *, size=18, weight="regular",
             color: RGBColor | None = None, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, letter_spacing: float | None = None,
             font: str = "Montserrat", line_spacing: float | None = None,
             upper: bool = False, brand: Brand | None = None):
    """Place a single-run text frame. Returns (textbox, frame, para, run).

    color defaults to brand.text_dark_primary (white).
    """
    if color is None:
        if brand is None:
            brand = load_brand()
        color = brand.text_dark_primary

    _, bold = WEIGHT_MAP[weight]
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing

    r = p.add_run()
    r.text = text.upper() if upper else text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color

    # letter-spacing — set via XML "spc" attr in 1/100 pt
    if letter_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))

    return tb, tf, p, r


def add_rect(slide, x, y, w, h, fill: RGBColor, *,
             line: RGBColor | None = None, line_w: float = 0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.shadow.inherit = False
    if line is not None:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_w or 0.75)
    else:
        rect.line.fill.background()
    return rect


def add_line(slide, x1, y1, x2, y2, color: RGBColor, weight: float = 1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ── signature motifs ──────────────────────────────────────────────────
def corner_brackets(slide, brand: Brand | None = None, *,
                    color: RGBColor | None = None, weight: float = 1.5,
                    size=Inches(0.35), margin=Inches(0.35)):
    """CAD-viewfinder corner brackets — Parspec primary motif.

    Applied to every slide as the persistent brand frame.
    """
    if brand is None:
        brand = load_brand()
    if color is None:
        color = brand.cta  # brand orange

    s, m = size, margin
    # top-left
    add_line(slide, m, m, m + s, m, color, weight)
    add_line(slide, m, m, m, m + s, color, weight)
    # top-right
    add_line(slide, SLIDE_W - m - s, m, SLIDE_W - m, m, color, weight)
    add_line(slide, SLIDE_W - m, m, SLIDE_W - m, m + s, color, weight)
    # bottom-left
    add_line(slide, m, SLIDE_H - m, m + s, SLIDE_H - m, color, weight)
    add_line(slide, m, SLIDE_H - m - s, m, SLIDE_H - m, color, weight)
    # bottom-right
    add_line(slide, SLIDE_W - m - s, SLIDE_H - m, SLIDE_W - m, SLIDE_H - m,
             color, weight)
    add_line(slide, SLIDE_W - m, SLIDE_H - m - s, SLIDE_W - m, SLIDE_H - m,
             color, weight)


def quarter_circle_accent(slide, x, y, *, radius=Inches(1.5),
                          color: RGBColor, weight: float = 2.0):
    """Quarter-circle arc — secondary CAD motif. Use sparingly."""
    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, x, y, radius, radius)
    arc.fill.background()
    arc.line.color.rgb = color
    arc.line.width = Pt(weight)
    arc.shadow.inherit = False
    return arc


def divider_rule(slide, x, y, w, *, color: RGBColor, weight: float = 2.0):
    """Short horizontal accent rule under section eyebrows."""
    add_line(slide, x, y, x + w, y, color, weight)


def vector_motif(slide, x, y, *, size=Inches(0.32), gap=Inches(0.04),
                 highlight: int | None = None,
                 color: RGBColor, hl_color: RGBColor):
    """2×2 grid of small squares — the recurring four-vector motif.

    highlight: None or int 0..3 (TL, TR, BL, BR) to color one square.
    """
    cells = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for idx, (cx, cy) in enumerate(cells):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + cx * (size + gap),
            y + cy * (size + gap),
            size, size,
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = (
            hl_color if highlight is not None and idx == highlight else color
        )
        rect.line.fill.background()
        rect.shadow.inherit = False
    return 2 * size + gap


# ── footer ────────────────────────────────────────────────────────────
def footer(slide, label: str, *, brand: Brand, page: int | None = None,
           total: int | None = None):
    """Bottom footer band: left label, right page indicator."""
    y = SLIDE_H - Inches(0.55)
    add_text(slide, Inches(0.6), y, Inches(6), Inches(0.3),
             label, size=8, weight="semibold",
             color=brand.text_dark_tertiary,
             letter_spacing=2.5, upper=True, brand=brand)
    if page is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.8), y, Inches(1.2), Inches(0.3),
                 f"{page:02d} / {total:02d}", size=8, weight="semibold",
                 color=brand.text_dark_tertiary,
                 align=PP_ALIGN.RIGHT, letter_spacing=2.5, brand=brand)
